"""
PPT Parser Service - Extract text and metadata from PowerPoint files
Uses python-pptx for parsing
"""

import io
import os
import textwrap
from pathlib import Path
from typing import Any

from common.error_handling.result import Result
from common.monitoring.logger import get_logger

logger = get_logger(__name__)

# Try to import python-pptx
try:
    from pptx import Presentation as PptxPresentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    PptxPresentation = None
    logger.warning("python-pptx not installed. PPT parsing will be limited.")

try:
    from PIL import Image, ImageDraw, ImageFont

    PILLOW_AVAILABLE = True
except ImportError:
    PILLOW_AVAILABLE = False
    Image = None
    ImageDraw = None
    ImageFont = None
    logger.warning("Pillow not installed. PPT thumbnail generation will be limited.")


class PPTParserService:
    """Service for parsing PowerPoint presentations"""

    def __init__(self):
        self.supported_formats = {".pptx", ".ppt"}

    async def parse_presentation(
        self, file_content: bytes, filename: str
    ) -> Result[dict[str, Any]]:
        """
        Parse a PPT file and extract text content from each slide

        Args:
            file_content: Raw bytes of the PPT file
            filename: Original filename for format detection

        Returns:
            Result containing parsed data with pages list
        """
        if not PPTX_AVAILABLE or PptxPresentation is None:
            return Result.fail("[PPTX_NOT_AVAILABLE] python-pptx library not installed")

        # Check file format
        file_ext = os.path.splitext(filename.lower())[1]
        if file_ext not in self.supported_formats:
            return Result.fail(f"[UNSUPPORTED_FORMAT] Format {file_ext} not supported")

        try:
            # Load presentation from bytes
            with io.BytesIO(file_content) as stream:
                prs = PptxPresentation(stream)

                pages = []
                total_slides = len(prs.slides)

                for idx, slide in enumerate(prs.slides, start=1):
                    page_data = self._extract_slide_content(slide, idx)
                    pages.append(page_data)

                result = {
                    "total_pages": total_slides,
                    "pages": pages,
                    "title": self._extract_title(prs),
                }

            logger.info(f"Parsed presentation: {total_slides} slides")
            return Result.ok(result)

        except Exception as e:
            logger.error(f"Failed to parse PPT: {str(e)}")
            return Result.fail(f"[PARSE_ERROR] {str(e)}")

    def _extract_slide_content(self, slide, page_number: int) -> dict[str, Any]:
        """Extract text content from a single slide"""
        texts = []

        # Extract text from all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

        # Join all text with newlines
        full_text = "\n".join(texts)

        # Extract notes if available
        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide:
            notes_text_frame = slide.notes_slide.notes_text_frame
            if notes_text_frame and notes_text_frame.text:
                notes_text = notes_text_frame.text.strip()

        return {
            "page_number": page_number,
            "extracted_text": full_text,
            "notes": notes_text,
            "text_length": len(full_text),
        }

    def _extract_title(self, prs: Any) -> str:
        """Extract presentation title from first slide or properties"""
        # Try to get from core properties
        if prs.core_properties and prs.core_properties.title:
            return prs.core_properties.title

        # Try to get from first slide title
        if len(prs.slides) > 0:
            first_slide = prs.slides[0]
            for shape in first_slide.shapes:
                if shape.is_placeholder:
                    if hasattr(shape, "text") and shape.text.strip():
                        return shape.text.strip()[:200]  # Limit length

        return "Untitled Presentation"

    async def generate_thumbnail(
        self,
        file_content: bytes,
        page_number: int = 1,
        output_dir: str = "./data/ppts/thumbnails",
    ) -> Result[str]:
        if (
            not PILLOW_AVAILABLE
            or Image is None
            or ImageDraw is None
            or ImageFont is None
        ):
            return Result.fail(
                "[PILLOW_NOT_AVAILABLE] Pillow is required for thumbnail generation"
            )

        if page_number < 1:
            return Result.fail("[INVALID_PAGE_NUMBER] page_number must be >= 1")

        slide_title = ""
        slide_text = ""

        if PPTX_AVAILABLE and PptxPresentation is not None:
            try:
                with io.BytesIO(file_content) as stream:
                    presentation = PptxPresentation(stream)
                    slide_count = len(presentation.slides)
                    if 1 <= page_number <= slide_count:
                        slide = presentation.slides[page_number - 1]
                        extracted = self._extract_slide_content(slide, page_number)
                        slide_text = extracted.get("extracted_text", "")
                    slide_title = self._extract_title(presentation)
            except Exception as exc:
                logger.warning(f"Thumbnail text extraction fallback: {exc}")

        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)
        thumbnail_path = output_path / f"page-{page_number}.png"

        image = Image.new("RGB", (1280, 720), "#f8fafc")
        draw = ImageDraw.Draw(image)

        try:
            title_font = ImageFont.truetype("Arial.ttf", 40)
            body_font = ImageFont.truetype("Arial.ttf", 28)
            footer_font = ImageFont.truetype("Arial.ttf", 24)
        except Exception:
            title_font = ImageFont.load_default()
            body_font = ImageFont.load_default()
            footer_font = ImageFont.load_default()

        draw.rectangle((0, 0, 1280, 120), fill="#1e293b")
        draw.text((36, 36), f"Slide {page_number}", fill="#f8fafc", font=title_font)

        if slide_title:
            draw.text((36, 140), slide_title[:80], fill="#0f172a", font=title_font)

        normalized_text = " ".join((slide_text or "").split())
        wrapped = textwrap.wrap(normalized_text, width=48)
        if not wrapped:
            wrapped = ["No text extracted from this slide."]

        y = 220
        for line in wrapped[:10]:
            draw.text((36, y), line, fill="#334155", font=body_font)
            y += 46

        draw.rectangle((0, 660, 1280, 720), fill="#e2e8f0")
        draw.text(
            (36, 676), "AI Presentation Practice", fill="#475569", font=footer_font
        )

        image.save(thumbnail_path, format="PNG", optimize=True)
        return Result.ok(str(thumbnail_path))


class PageContextService:
    """Service for managing page context during presentation sessions"""

    def __init__(self):
        self._session_pages: dict[str, dict[str, Any]] = {}

    def set_current_page(
        self, session_id: str, page_number: int, context: dict[str, Any]
    ) -> None:
        """Set current page context for a session"""
        if session_id not in self._session_pages:
            self._session_pages[session_id] = {}

        self._session_pages[session_id] = {
            "current_page": page_number,
            "total_pages": context.get("total_pages", page_number),
            "page_content": context.get("page_content", ""),
            "required_points": context.get("required_points", []),
            "forbidden_words": context.get("forbidden_words", []),
        }

        logger.debug(
            f"Updated page context for session {session_id}: page {page_number}"
        )

    def get_current_page(self, session_id: str) -> dict[str, Any] | None:
        """Get current page context for a session"""
        return self._session_pages.get(session_id)

    def clear_session(self, session_id: str) -> None:
        """Clear session context when session ends"""
        if session_id in self._session_pages:
            del self._session_pages[session_id]
            logger.debug(f"Cleared page context for session {session_id}")

    def is_valid_page(self, session_id: str, page_number: int) -> bool:
        """Check if page number is valid for session"""
        context = self._session_pages.get(session_id)
        if not context:
            return False

        total_pages = context.get("total_pages", 0)
        return 1 <= page_number <= total_pages


# Singleton instances
_ppt_parser: PPTParserService | None = None
_page_context: PageContextService | None = None


def get_ppt_parser() -> PPTParserService:
    """Get singleton PPT parser instance"""
    global _ppt_parser
    if _ppt_parser is None:
        _ppt_parser = PPTParserService()
    return _ppt_parser


def get_page_context_service() -> PageContextService:
    """Get singleton page context service instance"""
    global _page_context
    if _page_context is None:
        _page_context = PageContextService()
    return _page_context
