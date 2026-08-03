"""Controlled upload validation and preview production for learning sources.

This module deliberately owns only the Learning source pipeline.  It does not
broaden the legacy knowledge-base upload contract.
"""

from __future__ import annotations

import asyncio
import hashlib
import io
import json
import os
import re
import shutil
import stat
import subprocess
import tempfile
import unicodedata
import zipfile
from dataclasses import dataclass
from pathlib import Path, PurePosixPath
from typing import Any, Literal, cast
from xml.etree import ElementTree

from fastapi import UploadFile

from common.knowledge.processor import DocumentProcessor
from common.storage import DocumentStorageService

SourceContentKind = Literal[
    "document",
    "slide_deck",
    "demo_video",
    "external_demo",
    "script",
    "example_audio",
    "attachment",
]
SourceFileType = Literal[
    "pdf",
    "docx",
    "txt",
    "md",
    "xlsx",
    "xls",
    "pptx",
    "mp4",
    "webm",
    "mp3",
    "wav",
    "m4a",
]

DOCUMENT_FILE_TYPES = frozenset({"pdf", "docx", "txt", "md", "xlsx", "xls"})
SLIDE_FILE_TYPES = frozenset({"pptx"})
VIDEO_FILE_TYPES = frozenset({"mp4", "webm"})
AUDIO_FILE_TYPES = frozenset({"mp3", "wav", "m4a"})
SUPPORTED_SOURCE_FILE_TYPES = frozenset(
    DOCUMENT_FILE_TYPES | SLIDE_FILE_TYPES | VIDEO_FILE_TYPES | AUDIO_FILE_TYPES
)
RECOGNIZED_LEGACY_FILE_TYPES = frozenset({"ppt"})
PREVIEW_VERSION = "learning-preview-v1"

_TRUSTED_MIME_TYPES: dict[str, str] = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "txt": "text/plain",
    "md": "text/markdown",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "xls": "application/vnd.ms-excel",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "mp4": "video/mp4",
    "webm": "video/webm",
    "mp3": "audio/mpeg",
    "wav": "audio/wav",
    "m4a": "audio/mp4",
}
_KIND_FILE_TYPES: dict[str, frozenset[str]] = {
    "document": DOCUMENT_FILE_TYPES,
    "slide_deck": SLIDE_FILE_TYPES,
    "demo_video": VIDEO_FILE_TYPES,
    "example_audio": AUDIO_FILE_TYPES,
    "attachment": SUPPORTED_SOURCE_FILE_TYPES,
}
_SAFE_FILENAME_RE = re.compile(r"[^A-Za-z0-9._\-\u4e00-\u9fff]+")


class SourceUploadError(Exception):
    def __init__(self, code: str, message: str, *, status_code: int = 400) -> None:
        self.code = code
        self.message = message
        self.status_code = status_code
        super().__init__(message)


@dataclass(frozen=True, slots=True)
class SourceUploadPolicy:
    document_max_bytes: int
    media_max_bytes: int
    attachment_max_bytes: int
    zip_max_entries: int
    zip_max_uncompressed_bytes: int
    zip_max_ratio: int
    video_codecs: frozenset[str]
    audio_codecs: frozenset[str]

    @classmethod
    def from_env(cls) -> SourceUploadPolicy:
        return cls(
            document_max_bytes=_positive_env(
                "LEARNING_SOURCE_DOCUMENT_MAX_BYTES", 100 * 1024 * 1024
            ),
            media_max_bytes=_positive_env(
                "LEARNING_SOURCE_MEDIA_MAX_BYTES", 1024 * 1024 * 1024
            ),
            attachment_max_bytes=_positive_env(
                "LEARNING_SOURCE_ATTACHMENT_MAX_BYTES", 250 * 1024 * 1024
            ),
            zip_max_entries=_positive_env("LEARNING_SOURCE_ZIP_MAX_ENTRIES", 10_000),
            zip_max_uncompressed_bytes=_positive_env(
                "LEARNING_SOURCE_ZIP_MAX_UNCOMPRESSED_BYTES", 512 * 1024 * 1024
            ),
            zip_max_ratio=_positive_env("LEARNING_SOURCE_ZIP_MAX_RATIO", 100),
            video_codecs=_codec_env(
                "LEARNING_SOURCE_VIDEO_CODECS",
                frozenset({"h264", "vp8", "vp9", "av1"}),
                frozenset({"h264", "vp8", "vp9", "av1"}),
            ),
            audio_codecs=_codec_env(
                "LEARNING_SOURCE_AUDIO_CODECS",
                frozenset({"aac", "mp3", "opus", "vorbis", "pcm_s16le", "pcm_s24le"}),
                frozenset(
                    {"aac", "mp3", "opus", "vorbis", "pcm_s16le", "pcm_s24le"}
                ),
            ),
        )

    def max_bytes(self, content_kind: SourceContentKind) -> int:
        if content_kind in {"demo_video", "example_audio"}:
            return self.media_max_bytes
        if content_kind == "attachment":
            return self.attachment_max_bytes
        return self.document_max_bytes


@dataclass(frozen=True, slots=True)
class StagedSourceUpload:
    path: Path
    original_filename: str
    file_type: SourceFileType
    trusted_mime_type: str
    file_size_bytes: int
    file_hash: str

    def discard(self) -> None:
        self.path.unlink(missing_ok=True)


@dataclass(frozen=True, slots=True)
class SourceProcessingResult:
    processing_state: Literal["partial", "ready", "failed"]
    chunk_count: int
    artifact_available: bool
    manifest: dict[str, Any]
    anchors: tuple[dict[str, Any], ...]
    page_count: int | None = None
    duration_ms: int | None = None
    error_code: str | None = None
    error_message: str | None = None


async def stage_source_upload(
    upload: UploadFile,
    *,
    content_kind: SourceContentKind,
    storage: DocumentStorageService,
    policy: SourceUploadPolicy | None = None,
) -> StagedSourceUpload:
    """Stream to a tenant-neutral staging file before any database transaction."""

    policy = policy or SourceUploadPolicy.from_env()
    filename, extension = _safe_filename_and_extension(upload.filename or "")
    if extension == "ppt":
        raise SourceUploadError(
            "source_ppt_conversion_required",
            "当前环境不能可信转换旧版 .ppt，请先另存为 .pptx 后重新上传。",
            status_code=422,
        )
    if extension not in SUPPORTED_SOURCE_FILE_TYPES:
        raise SourceUploadError(
            "source_file_type_unsupported",
            "该文件格式暂不支持，请选择页面列出的受支持格式。",
            status_code=422,
        )
    allowed = _KIND_FILE_TYPES.get(content_kind)
    if allowed is None or extension not in allowed:
        raise SourceUploadError(
            "source_content_kind_mismatch",
            "材料类型与文件格式不匹配，请调整材料类型或重新选择文件。",
            status_code=422,
        )

    staging_dir = (storage.base_path / ".learning-source-staging").resolve()
    staging_dir.mkdir(parents=True, exist_ok=True)
    hasher = hashlib.sha256()
    size = 0
    max_bytes = policy.max_bytes(content_kind)
    descriptor, raw_path = tempfile.mkstemp(
        prefix="source-",
        suffix=f".{extension}.upload",
        dir=staging_dir,
    )
    path = Path(raw_path)
    try:
        with os.fdopen(descriptor, "wb") as stream:
            while True:
                chunk = await upload.read(1024 * 1024)
                if not chunk:
                    break
                size += len(chunk)
                if size > max_bytes:
                    raise SourceUploadError(
                        "source_file_too_large",
                        f"文件超过当前材料类型的 {max_bytes // (1024 * 1024)} MB 上限。",
                        status_code=413,
                    )
                hasher.update(chunk)
                stream.write(chunk)
            stream.flush()
            os.fsync(stream.fileno())
        if size == 0:
            raise SourceUploadError("source_file_empty", "不能上传空文件。", status_code=422)
        file_type = cast(SourceFileType, extension)
        validate_source_file(path, file_type=file_type, policy=policy)
        return StagedSourceUpload(
            path=path,
            original_filename=filename,
            file_type=file_type,
            trusted_mime_type=_TRUSTED_MIME_TYPES[file_type],
            file_size_bytes=size,
            file_hash=hasher.hexdigest(),
        )
    except Exception:
        path.unlink(missing_ok=True)
        raise
    finally:
        await upload.close()


def finalize_staged_source(staged: StagedSourceUpload, target: Path) -> bool:
    """Atomically move staged bytes into immutable controlled storage."""

    target = target.resolve()
    target.parent.mkdir(parents=True, exist_ok=True)
    if target.exists():
        if _sha256_path(target) != staged.file_hash:
            raise SourceUploadError(
                "source_artifact_conflict",
                "已保存材料与本次上传不一致，请联系管理员检查存储。",
                status_code=409,
            )
        staged.discard()
        return False
    os.replace(staged.path, target)
    return True


def validate_source_file(
    path: Path,
    *,
    file_type: SourceFileType,
    policy: SourceUploadPolicy | None = None,
) -> None:
    policy = policy or SourceUploadPolicy.from_env()
    with path.open("rb") as stream:
        head = stream.read(64)
    if file_type == "pdf" and not head.startswith(b"%PDF-"):
        _signature_mismatch("PDF")
    if file_type in {"docx", "xlsx", "pptx"}:
        _validate_ooxml(path, file_type=file_type, policy=policy)
    elif file_type == "xls" and not head.startswith(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"):
        _signature_mismatch("XLS")
    elif file_type in {"txt", "md"}:
        try:
            path.read_text(encoding="utf-8")
        except UnicodeDecodeError as exc:
            raise SourceUploadError(
                "source_text_encoding_invalid",
                "文本材料必须使用 UTF-8 编码。",
                status_code=422,
            ) from exc
    elif file_type == "wav" and not (
        head.startswith(b"RIFF") and head[8:12] == b"WAVE"
    ):
        _signature_mismatch("WAV")
    elif file_type == "mp3" and not (
        head.startswith(b"ID3")
        or (len(head) >= 2 and head[0] == 0xFF and head[1] & 0xE0 == 0xE0)
    ):
        _signature_mismatch("MP3")
    elif file_type in {"mp4", "m4a"} and not (
        len(head) >= 12 and head[4:8] == b"ftyp"
    ):
        _signature_mismatch(file_type.upper())
    elif file_type == "webm" and not head.startswith(b"\x1aE\xdf\xa3"):
        _signature_mismatch("WebM")


async def process_source_file(
    *,
    file_path: Path,
    file_type: SourceFileType,
    content_kind: SourceContentKind,
    storage: DocumentStorageService,
) -> SourceProcessingResult:
    if content_kind == "slide_deck":
        return await asyncio.to_thread(_render_pptx, file_path, storage)
    if content_kind in {"demo_video", "example_audio"}:
        return await asyncio.to_thread(
            _probe_media,
            file_path,
            file_type,
            content_kind,
            storage,
        )
    if content_kind == "attachment":
        attachment_hash = _sha256_path(file_path)
        manifest = {
            "version": PREVIEW_VERSION,
            "kind": "attachment",
            "download_available": True,
        }
        artifact_available = storage.save_parse_artifact(
            file_path, {"artifact_version": PREVIEW_VERSION, "manifest": manifest}
        ) is not None
        return SourceProcessingResult(
            processing_state="ready" if artifact_available else "failed",
            chunk_count=0,
            artifact_available=artifact_available,
            manifest=manifest if artifact_available else {},
            anchors=(
                {
                    "anchor_key": "full-attachment",
                    "label": "完整附件",
                    "locator": {
                        "type": "paragraph",
                        "paragraph_id": "full-attachment",
                        "start_offset": 0,
                        "end_offset": 1,
                    },
                    "excerpt_hash": attachment_hash,
                },
            )
            if artifact_available
            else (),
            error_code=None if artifact_available else "source_preview_store_failed",
            error_message=None if artifact_available else "附件预览信息保存失败，可稍后重试。",
        )

    parser_result = await DocumentProcessor().parse_document_artifact(
        file_path=str(file_path),
        file_type=file_type,
    )
    if str(parser_result.get("status")) != "ready":
        return SourceProcessingResult(
            processing_state="failed",
            chunk_count=0,
            artifact_available=False,
            manifest={},
            anchors=(),
            error_code=_document_failure_code(parser_result),
            error_message="材料内容没有成功解析，原文件已保留，可修正后重试。",
        )
    artifact = storage.load_parse_artifact(file_path) or {}
    chunks = artifact.get("chunks") if isinstance(artifact.get("chunks"), list) else []
    sections: list[dict[str, Any]] = []
    anchors: list[dict[str, Any]] = []
    for index, raw_chunk in enumerate(chunks, start=1):
        if not isinstance(raw_chunk, dict):
            continue
        content = str(raw_chunk.get("content") or "").strip()
        metadata = raw_chunk.get("metadata")
        metadata = metadata if isinstance(metadata, dict) else {}
        page = metadata.get("page")
        start_char = max(int(metadata.get("start_char") or 0), 0)
        end_char = max(int(metadata.get("end_char") or len(content)), start_char + 1)
        locator = (
            {"type": "page", "page": int(page), "start_offset": 0, "end_offset": max(len(content), 1)}
            if isinstance(page, int) and page > 0
            else {
                "type": "paragraph",
                "paragraph_id": f"chunk-{index}",
                "start_offset": start_char,
                "end_offset": end_char,
            }
        )
        sections.append({"index": index, "text": content, "locator": locator})
        anchors.append(
            {
                "anchor_key": f"chunk-{index}",
                "label": f"材料段落 {index}",
                "locator": locator,
                "excerpt_hash": hashlib.sha256(content.encode("utf-8")).hexdigest(),
            }
        )
    manifest = {
        "version": PREVIEW_VERSION,
        "kind": "document",
        "sections": sections,
    }
    return SourceProcessingResult(
        processing_state="ready",
        chunk_count=len(sections),
        artifact_available=True,
        manifest=manifest,
        anchors=tuple(anchors),
        page_count=_positive_int((artifact.get("metrics") or {}).get("page_count")),
    )


def preview_root(file_path: Path) -> Path:
    return file_path.with_name(f"{file_path.name}.preview")


def _render_pptx(
    file_path: Path,
    storage: DocumentStorageService,
) -> SourceProcessingResult:
    """Render actual slide geometry/images/text; never generate text-card stand-ins."""

    try:
        from PIL import Image, ImageDraw, ImageFont
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return _failed("source_slide_renderer_unavailable", "PPTX 分页渲染组件不可用，请稍后重试。")

    root = preview_root(file_path)
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    pages: list[dict[str, Any]] = []
    anchors: list[dict[str, Any]] = []
    failed_pages: list[int] = []
    try:
        presentation = Presentation(str(file_path))
    except Exception:
        return _failed("source_slide_decode_failed", "PPTX 文件无法解码，请检查文件后重试。")
    slide_width = max(int(presentation.slide_width or 1), 1)
    slide_height = max(int(presentation.slide_height or 1), 1)
    width = 1280
    height = max(round(width * slide_height / slide_width), 1)
    font_path = _font_path()

    for page_number, slide in enumerate(presentation.slides, start=1):
        extracted_text: list[str] = []
        try:
            image = Image.new("RGB", (width, height), "white")
            draw = ImageDraw.Draw(image)
            for shape in slide.shapes:
                x = round(int(shape.left) * width / slide_width)
                y = round(int(shape.top) * height / slide_height)
                w = max(round(int(shape.width) * width / slide_width), 1)
                h = max(round(int(shape.height) * height / slide_height), 1)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    picture = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                    picture.thumbnail((w, h))
                    image.paste(picture, (x, y))
                elif getattr(shape, "has_text_frame", False):
                    text = str(getattr(shape, "text", "") or "").strip()
                    if not text:
                        continue
                    extracted_text.append(text)
                    fill = _shape_fill(shape)
                    if fill:
                        draw.rectangle((x, y, x + w, y + h), fill=fill)
                    size = max(min(round(h / max(text.count("\n") + 2, 2)), 36), 14)
                    font = ImageFont.truetype(font_path, size) if font_path else ImageFont.load_default()
                    draw.multiline_text(
                        (x + 6, y + 4),
                        _fit_text(text, max(round(w / max(size * 0.65, 1)), 8)),
                        fill="#111827",
                        font=font,
                        spacing=4,
                    )
                elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    draw.rectangle((x, y, x + w, y + h), outline="#64748b", width=2)
            output = root / f"page-{page_number}.png"
            image.save(output, format="PNG", optimize=True)
            text = "\n".join(extracted_text).strip()
            page_entry = {
                "page": page_number,
                "status": "ready",
                "preview_ref": output.name,
                "text": text,
            }
            pages.append(page_entry)
            anchors.append(
                {
                    "anchor_key": f"slide-{page_number}",
                    "label": f"第 {page_number} 页",
                    "locator": {
                        "type": "page",
                        "page": page_number,
                        "start_offset": 0,
                        "end_offset": max(len(text), 1),
                    },
                    "excerpt_hash": hashlib.sha256(text.encode("utf-8")).hexdigest(),
                }
            )
        except Exception:
            failed_pages.append(page_number)
            pages.append({"page": page_number, "status": "failed"})

    ready_pages = [item for item in pages if item.get("status") == "ready"]
    state: Literal["partial", "ready", "failed"]
    if ready_pages and not failed_pages:
        state = "ready"
    elif ready_pages:
        state = "partial"
    else:
        state = "failed"
    manifest = {
        "version": PREVIEW_VERSION,
        "kind": "slide_deck",
        "page_count": len(pages),
        "pages": pages,
        "missing_pages": failed_pages,
        "rendering": "controlled_raster",
    }
    artifact = {
        "artifact_version": PREVIEW_VERSION,
        "file_type": "pptx",
        "manifest": manifest,
        "chunks": [
            {
                "index": index,
                "content": item.get("text", ""),
                "metadata": {"page": item["page"]},
            }
            for index, item in enumerate(ready_pages)
        ],
    }
    artifact_available = storage.save_parse_artifact(file_path, artifact) is not None
    if not artifact_available:
        return _failed("source_preview_store_failed", "分页预览没有成功保存，可稍后重试。")
    return SourceProcessingResult(
        processing_state=state,
        chunk_count=len(ready_pages),
        artifact_available=True,
        manifest=manifest,
        anchors=tuple(anchors),
        page_count=len(pages) or None,
        error_code="source_slide_pages_partial" if state == "partial" else (
            "source_slide_render_failed" if state == "failed" else None
        ),
        error_message=(
            f"第 {', '.join(str(item) for item in failed_pages)} 页预览失败；成功页和原文件已保留。"
            if state == "partial"
            else "PPTX 没有可用的分页预览；原文件已保留。"
            if state == "failed"
            else None
        ),
    )


def _probe_media(
    file_path: Path,
    file_type: SourceFileType,
    content_kind: SourceContentKind,
    storage: DocumentStorageService,
) -> SourceProcessingResult:
    executable = shutil.which("ffprobe")
    if executable is None:
        return _failed(
            "source_media_probe_unavailable",
            "媒体解码服务暂不可用，原文件已保留，可稍后重试。",
        )
    try:
        completed = subprocess.run(
            [
                executable,
                "-v",
                "error",
                "-show_entries",
                "format=duration,format_name:stream=codec_type,codec_name",
                "-of",
                "json",
                str(file_path),
            ],
            check=False,
            capture_output=True,
            text=True,
            timeout=30,
        )
        if completed.returncode != 0:
            return _failed("source_media_decode_failed", "媒体文件无法解码，请更换文件后重试。")
        probe = json.loads(completed.stdout)
    except (OSError, subprocess.SubprocessError, json.JSONDecodeError):
        return _failed("source_media_probe_failed", "媒体文件探测失败，原文件已保留，可稍后重试。")
    streams = probe.get("streams") if isinstance(probe, dict) else None
    streams = streams if isinstance(streams, list) else []
    required_type = "video" if content_kind == "demo_video" else "audio"
    matching = [
        item
        for item in streams
        if isinstance(item, dict) and item.get("codec_type") == required_type
    ]
    try:
        policy = SourceUploadPolicy.from_env()
    except SourceUploadError:
        return _failed("source_media_policy_invalid", "媒体编码策略配置无效，请联系管理员。")
    allowed_codecs = (
        policy.video_codecs if required_type == "video" else policy.audio_codecs
    )
    if not matching or any(str(item.get("codec_name")) not in allowed_codecs for item in matching):
        return _failed("source_media_codec_unsupported", "媒体编码不受支持，请转换格式后重试。")
    duration_raw = (probe.get("format") or {}).get("duration")
    try:
        duration_ms = round(float(duration_raw) * 1000)
    except (TypeError, ValueError):
        duration_ms = 0
    if duration_ms <= 0:
        return _failed("source_media_duration_invalid", "媒体时长无法确认，请更换文件后重试。")
    root = preview_root(file_path)
    root.mkdir(parents=True, exist_ok=True)
    playback = root / f"playback.{file_type}"
    shutil.copyfile(file_path, playback)
    manifest = {
        "version": PREVIEW_VERSION,
        "kind": content_kind,
        "duration_ms": duration_ms,
        "playback_ref": playback.name,
        "codecs": [str(item.get("codec_name")) for item in matching],
    }
    artifact_available = storage.save_parse_artifact(
        file_path,
        {"artifact_version": PREVIEW_VERSION, "manifest": manifest},
    ) is not None
    if not artifact_available:
        return _failed("source_preview_store_failed", "媒体播放信息没有成功保存，可稍后重试。")
    return SourceProcessingResult(
        processing_state="ready",
        chunk_count=0,
        artifact_available=True,
        manifest=manifest,
        anchors=(
            {
                "anchor_key": "full-media",
                "label": "完整媒体",
                "locator": {"type": "time_range", "start_ms": 0, "end_ms": duration_ms},
                "excerpt_hash": hashlib.sha256(
                    f"0:{duration_ms}".encode()
                ).hexdigest(),
            },
        ),
        duration_ms=duration_ms,
    )


def _validate_ooxml(
    path: Path,
    *,
    file_type: SourceFileType,
    policy: SourceUploadPolicy,
) -> None:
    required = {
        "docx": "word/document.xml",
        "xlsx": "xl/workbook.xml",
        "pptx": "ppt/presentation.xml",
    }[file_type]
    required_content_type = {
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
    }[file_type]
    try:
        with zipfile.ZipFile(path) as package:
            infos = package.infolist()
            if len(infos) > policy.zip_max_entries:
                raise SourceUploadError(
                    "source_zip_too_many_entries",
                    "Office 文件包含过多内部条目，已拒绝处理。",
                    status_code=422,
                )
            names = {item.filename for item in infos}
            if required not in names or "[Content_Types].xml" not in names:
                _signature_mismatch(file_type.upper())
            uncompressed = 0
            compressed = 0
            for item in infos:
                pure = PurePosixPath(item.filename)
                if pure.is_absolute() or ".." in pure.parts or "\\" in item.filename:
                    raise SourceUploadError(
                        "source_zip_path_invalid",
                        "Office 文件包含不安全的内部路径，已拒绝处理。",
                        status_code=422,
                    )
                mode = item.external_attr >> 16
                if mode and stat.S_ISLNK(mode):
                    raise SourceUploadError(
                        "source_zip_link_invalid",
                        "Office 文件包含不安全的内部链接，已拒绝处理。",
                        status_code=422,
                    )
                if item.flag_bits & 0x1:
                    raise SourceUploadError(
                        "source_zip_encrypted",
                        "暂不支持加密 Office 文件，请解除密码后重试。",
                        status_code=422,
                    )
                uncompressed += item.file_size
                compressed += item.compress_size
            if uncompressed > policy.zip_max_uncompressed_bytes or (
                compressed > 0 and uncompressed / compressed > policy.zip_max_ratio
            ):
                raise SourceUploadError(
                    "source_zip_bomb_detected",
                    "Office 文件解压规模异常，已拒绝处理。",
                    status_code=422,
                )
            if package.testzip() is not None:
                raise SourceUploadError(
                    "source_zip_crc_invalid",
                    "Office 文件内部校验失败，请重新导出后上传。",
                    status_code=422,
                )
            try:
                content_types_root = ElementTree.fromstring(
                    package.read("[Content_Types].xml")
                )
            except (ElementTree.ParseError, KeyError) as exc:
                raise SourceUploadError(
                    "source_ooxml_content_types_invalid",
                    "Office 文件的内容类型清单无效，已拒绝处理。",
                    status_code=422,
                ) from exc
            declared = {
                (
                    str(item.attrib.get("PartName") or "").lstrip("/"),
                    str(item.attrib.get("ContentType") or ""),
                )
                for item in content_types_root
                if item.tag.rsplit("}", 1)[-1] == "Override"
            }
            if (required, required_content_type) not in declared:
                raise SourceUploadError(
                    "source_ooxml_content_type_mismatch",
                    f"文件不是声明完整的 {file_type.upper()} 文档。",
                    status_code=422,
                )
    except zipfile.BadZipFile as exc:
        raise SourceUploadError(
            "source_ooxml_invalid",
            f"文件内容不是有效的 {file_type.upper()}。",
            status_code=422,
        ) from exc


def _safe_filename_and_extension(filename: str) -> tuple[str, str]:
    normalized = unicodedata.normalize("NFKC", Path(filename).name).strip()
    normalized = _SAFE_FILENAME_RE.sub("_", normalized)[:255]
    if not normalized or "." not in normalized:
        raise SourceUploadError(
            "source_filename_invalid",
            "文件名必须包含受支持的扩展名。",
            status_code=422,
        )
    extension = normalized.rsplit(".", 1)[-1].lower()
    return normalized, extension


def _positive_env(name: str, default: int) -> int:
    try:
        value = int(os.getenv(name, str(default)))
    except ValueError:
        return default
    return value if value > 0 else default


def _codec_env(
    name: str,
    default: frozenset[str],
    supported: frozenset[str],
) -> frozenset[str]:
    raw = os.getenv(name)
    if raw is None:
        return default
    configured = frozenset(item.strip().lower() for item in raw.split(",") if item.strip())
    if not configured or not configured.issubset(supported):
        raise SourceUploadError(
            "source_policy_invalid",
            f"{name} 包含未支持的编码或为空。",
            status_code=503,
        )
    return configured


def _sha256_path(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _signature_mismatch(label: str) -> None:
    raise SourceUploadError(
        "source_signature_mismatch",
        f"文件内容不是有效的 {label}，请检查扩展名和文件内容。",
        status_code=422,
    )


def _document_failure_code(result: dict[str, Any]) -> str:
    message = str(result.get("error_message") or "")
    if "PARSE_ARTIFACT_SAVE_FAILED" in message:
        return "source_preview_store_failed"
    if "PARSE_EMPTY_STRUCTURED_DOC" in message:
        return "source_document_content_empty"
    return "source_document_parse_failed"


def _failed(code: str, message: str) -> SourceProcessingResult:
    return SourceProcessingResult(
        processing_state="failed",
        chunk_count=0,
        artifact_available=False,
        manifest={},
        anchors=(),
        error_code=code,
        error_message=message,
    )


def _positive_int(value: object) -> int | None:
    try:
        number = int(cast(Any, value))
    except (TypeError, ValueError):
        return None
    return number if number > 0 else None


def _shape_fill(shape: Any) -> str | None:
    try:
        rgb = shape.fill.fore_color.rgb
        return f"#{rgb}" if rgb is not None else None
    except (AttributeError, TypeError, ValueError):
        return None


def _fit_text(value: str, width: int) -> str:
    lines: list[str] = []
    for source_line in value.splitlines() or [value]:
        if not source_line:
            lines.append("")
            continue
        lines.extend(source_line[index : index + width] for index in range(0, len(source_line), width))
    return "\n".join(lines[:20])


def _font_path() -> str | None:
    candidates = (
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    )
    return next((item for item in candidates if Path(item).is_file()), None)


__all__ = [
    "AUDIO_FILE_TYPES",
    "DOCUMENT_FILE_TYPES",
    "PREVIEW_VERSION",
    "RECOGNIZED_LEGACY_FILE_TYPES",
    "SLIDE_FILE_TYPES",
    "SUPPORTED_SOURCE_FILE_TYPES",
    "SourceContentKind",
    "SourceFileType",
    "SourceProcessingResult",
    "SourceUploadError",
    "SourceUploadPolicy",
    "StagedSourceUpload",
    "VIDEO_FILE_TYPES",
    "finalize_staged_source",
    "preview_root",
    "process_source_file",
    "stage_source_upload",
    "validate_source_file",
]
