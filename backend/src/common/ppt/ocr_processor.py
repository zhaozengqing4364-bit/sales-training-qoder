"""
PPT OCR Processing Service - Extracts text from PPT slides

Implements Constitution Principles:
- I. NO ERROR POPUPS - Graceful degradation on OCR failure
- V. Cost control - Uses free OCR solutions
"""

import logging
import uuid
from dataclasses import dataclass

# For PPT text extraction (python-pptx doesn't require OCR for text-based PPTs)
# For image-based slides, would use Tesseract OCR
from common.error_handling.result import Result

logger = logging.getLogger(__name__)


@dataclass
class PPTPage:
    """Represents a single PPT page with extracted content"""
    page_number: int
    title: str
    content: str
    image_count: int


@dataclass
class PPTExtraction:
    """Result of PPT text extraction"""
    presentation_id: uuid.UUID
    filename: str
    total_pages: int
    pages: list[PPTPage]
    has_images: bool


class OCRProcessor:
    """
    Extracts text from PPT files

    Strategy:
    1. For text-based PPTs: Use python-pptx (fast, accurate)
    2. For image-based slides: Use Tesseract OCR (slower, fallback)
    3. On failure: Return empty extraction (graceful degradation)
    """

    def __init__(self):
        self.temp_dir = "/tmp/ppt_processing"

    async def extract_text(
        self,
        file_path: str,
        presentation_id: uuid.UUID,
        filename: str
    ) -> Result[PPTExtraction]:
        """
        Extract text from PPT file

        Returns: PPTExtraction or Result.fail
        """
        try:
            # Import here to avoid dependency if not used
            from pptx import Presentation

            # Load PPT
            prs = Presentation(file_path)

            pages = []
            has_images = False

            for slide_number, slide in enumerate(prs.slides, start=1):
                # Extract text from slide
                title, content, image_count = self._extract_slide_content(slide)

                if image_count > 0:
                    has_images = True

                page = PPTPage(
                    page_number=slide_number,
                    title=title,
                    content=content,
                    image_count=image_count
                )

                pages.append(page)

            extraction = PPTExtraction(
                presentation_id=presentation_id,
                filename=filename,
                total_pages=len(prs.slides),
                pages=pages,
                has_images=has_images
            )

            logger.info(
                "PPT text extraction complete",
                extra={
                    "presentation_id": str(presentation_id),
                    "pages": extraction.total_pages,
                    "has_images": has_images,
                }
            )

            return Result(value=extraction)

        except ImportError:
            logger.error("python-pptx not installed")
            return Result.fail(fallback="[PPTX_NOT_INSTALLED]")
        except (RuntimeError, ValueError, OSError) as e:
            logger.error(
                "Failed to extract text from PPT",
                extra={"presentation_id": str(presentation_id), "error": str(e)},
                exc_info=True
            )
            return Result.fail(fallback="[OCR_FAILED]")

    def _extract_slide_content(self, slide) -> tuple[str, str, int]:
        """
        Extract title, content, and image count from a slide

        Returns: (title, content, image_count)
        """
        title = ""
        content_lines = []
        image_count = 0

        # Extract text from shapes
        for shape in slide.shapes:
            # Check for images
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_count += 1
                continue

            # Extract text from text frames
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()

                if not text:
                    continue

                # First text frame is usually the title
                if not title and len(text) < 100:
                    title = text
                else:
                    content_lines.append(text)

        content = "\n".join(content_lines)

        # Fallback title if none found
        if not title:
            title = f"Slide {slide.slide_index + 1}"

        return title, content, image_count

    async def extract_with_ocr(
        self,
        image_path: str
    ) -> Result[str]:
        """
        Extract text from image using OCR (fallback for image-based slides)

        Uses Tesseract OCR (free, local)

        Returns: Extracted text or Result.fail
        """
        try:
            import pytesseract
            from PIL import Image

            # Open image
            image = Image.open(image_path)

            # Extract text
            text = pytesseract.image_to_string(image, lang='chi_sim+eng')

            logger.info(
                "OCR extraction complete",
                extra={"text_length": len(text)}
            )

            return Result(value=text.strip())

        except ImportError:
            logger.error("Tesseract OCR not installed")
            return Result.fail(fallback="[OCR_NOT_INSTALLED]")
        except (RuntimeError, ValueError, OSError) as e:
            logger.error(
                "OCR extraction failed",
                extra={"error": str(e)},
                exc_info=True
            )
            return Result.fail(fallback="[OCR_FAILED]")


# Singleton instance
ocr_processor = OCRProcessor()
